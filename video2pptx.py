import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches
import os

def extract_slide_frames(input_video_path, output_pptx_path="output.pptx", cleanup=True, debug=False, matcher=cv2.TM_CCOEFF_NORMED):

    # Define thresholds for detecting slide changes and additional content
    SLIDE_THRESHOLD   = 0.2
    CONTENT_THRESHOLD = 0.98
    width, height = Inches(10), Inches(5.625)  # Width is 10 inches, Height is 5.625 inches (10/16 * 9)
    SECONDS_BETWEEN_FRAMES = 5

    # Open the video file
    cap = cv2.VideoCapture(input_video_path)

    # Initialize variables to store previous frame and slide number
    prev_frame = None
    slide_number = 0
    frame_rate = int(cap.get(cv2.CAP_PROP_FPS))  # Get the frame rate of the video

    # Define the region of interest (ROI) 
    roi_x, roi_y, roi_width, roi_height = 0, 0, int(cap.get(cv2.CAP_PROP_FRAME_WIDTH)), int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
    # Example: Exclude 10% of width from the right
    #roi_width -= int(roi_width * 0.1)  

    # Create a PowerPoint presentation object with widescreen aspect ratio (16:9)
    prs = Presentation()
    prs.slide_width = width
    prs.slide_height = height

    frame_count = 0
    slide = None  # Initialize slide variable outside of the loop

    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break

        frame_count += 1

        # Skip frames based on frame rate
        if frame_count % (frame_rate * SECONDS_BETWEEN_FRAMES) != 0:
            continue

        # Extract the region of interest (ROI) from the frame
        roi_frame = frame[roi_y:roi_y + roi_height, roi_x:roi_x + roi_width]

        # Convert the ROI frame to grayscale for comparison
        gray_frame = cv2.cvtColor(roi_frame, cv2.COLOR_BGR2GRAY)

        # Initialize prev_frame if it's the first frame
        if prev_frame is None:
            prev_frame = gray_frame
            continue

        # Calculate structural similarity index between the current frame and the previous frame, cv2.TM_CCOEFF_NORMED worked best for my tests
        similarity_index = cv2.matchTemplate(prev_frame, gray_frame, matcher)

        # Take the absolute value of the similarity index
        similarity_index = np.abs(similarity_index)

        # Check if the similarity index is below the slide change threshold
        if np.max(similarity_index) < SLIDE_THRESHOLD:
            # If the similarity index is below the slide change threshold, it indicates a slide transition
            if debug:
                print(f"Found new slide {slide_number} at minute {cap.get(cv2.CAP_PROP_POS_MSEC) / 60000:.2f}. Similarity: {np.max(similarity_index)}")
            # Add a slide to the presentation
            slide_layout = prs.slide_layouts[5]  # Use a blank layout for the slide
            slide = prs.slides.add_slide(slide_layout)

            # Save the original frame as an image since the pptx lib does not understand bytes and uses from_file
            cv2.imwrite(f"slide_{slide_number}.png", frame)

            # Add the image to the slide
            left = top = 0
            slide.shapes.add_picture(f"slide_{slide_number}.png", left, top, width=width, height=height)

            # Update the previous frame
            prev_frame = gray_frame

            # Increment the slide number
            slide_number += 1

        elif slide is not None and np.max(similarity_index) < CONTENT_THRESHOLD:
            # If the similarity index is below the content threshold, it indicates additional content on the current slide
            if debug:
                print(f"Updating slide {slide_number} at minute {cap.get(cv2.CAP_PROP_POS_MSEC) / 60000:.2f}. Similarity: {np.max(similarity_index)}")
            # Save the original frame as an image
            cv2.imwrite(f"slide_{slide_number}.png", frame)

            # Add the image to the current slide
            slide.shapes.add_picture(f"slide_{slide_number}.png", left, top, width=width, height=height)

    # Save the PowerPoint presentation
    prs.save(output_pptx_path)

    # Release the video capture 
    cap.release()
    if cleanup:
        # Remove the slide files, default True
        for slide_number in range(slide_number):
            slide_file = f"slide_{slide_number}.png"
            if os.path.exists(slide_file):
                os.remove(slide_file)

# Example usage:
output_pptx_path = "output.pptx"
video_path = "sd.mp4"
#matcher = cv2.TM_SQDIFF_NORMED  
extract_slide_frames(video_path, output_pptx_path, True, True)
